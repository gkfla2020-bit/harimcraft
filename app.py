"""
하림님 전용 코딩 챗봇 v2.0
- 다중 채팅방 + 자동 저장
- 스트리밍 응답
- 채팅 검색/내보내기
- 테마 설정
"""
import os, io, traceback, json, re
from datetime import datetime
from typing import List, Dict
from dotenv import load_dotenv
from fastapi import FastAPI, UploadFile, File, Form, Query
from fastapi.responses import HTMLResponse, JSONResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from anthropic import Anthropic, APIConnectionError, RateLimitError, APIStatusError
import PyPDF2
import csv
import zipfile
import docx  # python-docx
import openpyxl  # 엑셀
import pptx  # python-pptx
import httpx  # 웹 검색용
from urllib.parse import quote_plus

load_dotenv()
app = FastAPI()
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

client = Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"), timeout=180.0, max_retries=3)

DATA_DIR = "data"
CHATS_FILE = os.path.join(DATA_DIR, "chats.json")
SETTINGS_FILE = os.path.join(DATA_DIR, "settings.json")

os.makedirs(DATA_DIR, exist_ok=True)

chats: Dict[str, dict] = {}
settings: dict = {"theme": "dark", "fontSize": "medium"}

def load_data():
    global chats, settings
    if os.path.exists(CHATS_FILE):
        try:
            with open(CHATS_FILE, 'r', encoding='utf-8') as f:
                chats = json.load(f)
        except: chats = {}
    if os.path.exists(SETTINGS_FILE):
        try:
            with open(SETTINGS_FILE, 'r', encoding='utf-8') as f:
                settings = json.load(f)
        except: pass

def save_chats():
    with open(CHATS_FILE, 'w', encoding='utf-8') as f:
        json.dump(chats, f, ensure_ascii=False, indent=2)

def save_settings():
    with open(SETTINGS_FILE, 'w', encoding='utf-8') as f:
        json.dump(settings, f, ensure_ascii=False, indent=2)

load_data()

SYSTEM_PROMPT = """당신은 정하림님의 개인 AI 어시스턴트입니다.
당신은 Claude Opus 4 모델입니다 (2025년 5월 버전, Anthropic 최고 성능 모델).

[하림님 정보] 퀀트 연구자, 논문 리뷰, Python 주력

[응답 원칙]
1. 어려운 용어는 쉽게 풀어서 설명, 비유/예시 많이
2. Python 코드에는 한글 주석 꼼꼼히
3. 논문은 핵심 아이디어, 방법론, 결과 위주로 요약
4. 수학 수식은 LaTeX (인라인: $수식$, 블록: $$수식$$)
5. 퀀트는 실제 트레이딩/백테스팅 관점에서 설명
6. pandas, numpy, scipy 활용 코드 제공
7. 코드 블록에는 반드시 언어 명시 (```python 등)

[퀀트/논문 지원]
- 팩터 투자, 포트폴리오 최적화, 리스크 관리
- 백테스팅 코드는 실제 사용 가능한 수준
- 논문: Abstract, Methodology, Results, Conclusion 구조 요약
- 수식은 Python 코드로도 구현해서 보여주기
"""

CACHED_SYSTEM = [{"type": "text", "text": SYSTEM_PROMPT, "cache_control": {"type": "ephemeral"}}]

def extract_pdf_text(pdf_bytes):
    try:
        reader = PyPDF2.PdfReader(io.BytesIO(pdf_bytes))
        return "\n".join(p.extract_text() or "" for p in reader.pages).strip()
    except: return ""

def extract_docx_text(docx_bytes):
    """Word 문서에서 텍스트 추출"""
    try:
        doc = docx.Document(io.BytesIO(docx_bytes))
        return "\n".join(p.text for p in doc.paragraphs).strip()
    except: return ""

def extract_xlsx_text(xlsx_bytes):
    """엑셀에서 텍스트 추출"""
    try:
        wb = openpyxl.load_workbook(io.BytesIO(xlsx_bytes), data_only=True)
        text = []
        for sheet in wb.worksheets:
            text.append(f"[시트: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    text.append(row_text)
        return "\n".join(text)
    except: return ""

def extract_pptx_text(pptx_bytes):
    """파워포인트에서 텍스트 추출"""
    try:
        prs = pptx.Presentation(io.BytesIO(pptx_bytes))
        text = []
        for i, slide in enumerate(prs.slides, 1):
            text.append(f"[슬라이드 {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return "\n".join(text)
    except: return ""

def extract_csv_text(csv_bytes):
    """CSV에서 텍스트 추출"""
    try:
        content = csv_bytes.decode('utf-8-sig')
        return content[:50000]
    except:
        try:
            content = csv_bytes.decode('cp949')
            return content[:50000]
        except: return ""

def extract_text_file(file_bytes, filename):
    """일반 텍스트 파일 추출"""
    try:
        return file_bytes.decode('utf-8')[:50000]
    except:
        try:
            return file_bytes.decode('cp949')[:50000]
        except: return ""

def extract_file_content(file_bytes, filename):
    """파일 형식에 따라 텍스트 추출"""
    ext = filename.lower().split('.')[-1] if '.' in filename else ''
    
    if ext == 'pdf':
        return extract_pdf_text(file_bytes)
    elif ext == 'docx':
        return extract_docx_text(file_bytes)
    elif ext in ['xlsx', 'xls']:
        return extract_xlsx_text(file_bytes)
    elif ext == 'pptx':
        return extract_pptx_text(file_bytes)
    elif ext == 'csv':
        return extract_csv_text(file_bytes)
    elif ext in ['txt', 'md', 'py', 'js', 'ts', 'java', 'c', 'cpp', 'h', 'json', 'xml', 'html', 'css', 'sql', 'yaml', 'yml', 'ini', 'cfg', 'log', 'sh', 'bat']:
        return extract_text_file(file_bytes, filename)
    else:
        # 알 수 없는 형식은 텍스트로 시도
        return extract_text_file(file_bytes, filename)

def generate_title(message: str) -> str:
    """첫 메시지로 제목 생성"""
    msg = message.strip()
    if len(msg) > 40:
        return msg[:37] + "..."
    return msg or "새 채팅"

async def web_search(query: str, num_results: int = 5) -> str:
    """웹 검색 - 여러 방법 시도"""
    try:
        async with httpx.AsyncClient(timeout=15.0, follow_redirects=True) as client:
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
                "Accept-Language": "ko-KR,ko;q=0.9,en-US;q=0.8,en;q=0.7",
            }
            
            results = []
            
            # 방법 1: DuckDuckGo HTML 검색
            try:
                search_url = f"https://html.duckduckgo.com/html/?q={quote_plus(query)}"
                response = await client.get(search_url, headers=headers)
                
                if response.status_code == 200:
                    html = response.text
                    import re
                    
                    # 결과 블록 파싱
                    result_blocks = re.findall(
                        r'<a[^>]*class="result__a"[^>]*href="([^"]*)"[^>]*>([^<]*)</a>.*?'
                        r'<a[^>]*class="result__snippet"[^>]*>([^<]*)</a>',
                        html, re.DOTALL
                    )
                    
                    if not result_blocks:
                        # 다른 패턴 시도
                        titles = re.findall(r'class="result__a"[^>]*>([^<]+)</a>', html)
                        snippets = re.findall(r'class="result__snippet"[^>]*>([^<]+)', html)
                        urls = re.findall(r'class="result__url"[^>]*>([^<]+)', html)
                        
                        for i in range(min(len(titles), len(snippets), num_results)):
                            title = titles[i].strip() if i < len(titles) else ""
                            snippet = snippets[i].strip() if i < len(snippets) else ""
                            url = urls[i].strip() if i < len(urls) else ""
                            if title and snippet:
                                results.append(f"**{title}**\n{snippet}\n🔗 {url}")
                    else:
                        for url, title, snippet in result_blocks[:num_results]:
                            if title.strip() and snippet.strip():
                                results.append(f"**{title.strip()}**\n{snippet.strip()}")
            except Exception as e:
                print(f"DuckDuckGo HTML 검색 실패: {e}")
            
            # 방법 2: DuckDuckGo Instant Answer API (위키피디아 등)
            if len(results) < 2:
                try:
                    api_url = f"https://api.duckduckgo.com/?q={quote_plus(query)}&format=json&no_html=1&skip_disambig=1"
                    response = await client.get(api_url, headers=headers)
                    data = response.json()
                    
                    # Abstract
                    if data.get("Abstract"):
                        source = data.get("AbstractSource", "")
                        results.insert(0, f"📖 **{source}**\n{data['Abstract']}")
                    
                    # Answer
                    if data.get("Answer"):
                        results.insert(0, f"💡 **답변**\n{data['Answer']}")
                    
                    # Related Topics
                    for topic in data.get("RelatedTopics", [])[:3]:
                        if isinstance(topic, dict) and topic.get("Text"):
                            text = topic["Text"]
                            if text not in str(results):
                                results.append(f"• {text}")
                except Exception as e:
                    print(f"DuckDuckGo API 검색 실패: {e}")
            
            # 방법 3: Wikipedia API 직접 검색 (한국어)
            if len(results) < 2:
                try:
                    wiki_url = f"https://ko.wikipedia.org/api/rest_v1/page/summary/{quote_plus(query)}"
                    response = await client.get(wiki_url, headers=headers)
                    if response.status_code == 200:
                        data = response.json()
                        if data.get("extract"):
                            results.append(f"📚 **위키백과: {data.get('title', query)}**\n{data['extract']}")
                except:
                    pass
                
                # 영어 위키피디아도 시도
                try:
                    wiki_url = f"https://en.wikipedia.org/api/rest_v1/page/summary/{quote_plus(query)}"
                    response = await client.get(wiki_url, headers=headers)
                    if response.status_code == 200:
                        data = response.json()
                        if data.get("extract"):
                            results.append(f"📚 **Wikipedia: {data.get('title', query)}**\n{data['extract']}")
                except:
                    pass
            
            if results:
                return "\n\n---\n\n".join(results[:num_results])
            return ""
            
    except Exception as e:
        print(f"웹 검색 전체 실패: {e}")
        return ""

def should_search(message: str) -> tuple[bool, str]:
    """메시지에서 검색 필요 여부와 검색어 추출"""
    import re
    msg = message.strip()
    
    # 명시적 검색 요청
    explicit_patterns = [
        r'(?:검색|찾아|알아)[해줘봐\s]*[:\s]*(.+)',
        r'(.+?)(?:에 대해|에대해)?\s*(?:검색|찾아|알아)[줘봐]',
        r'(.+?)\s*(?:뭐야|뭔가요|무엇인가요|이 뭐야)\??',
        r'(.+?)\s*(?:알려줘|설명해줘|가르쳐줘)',
    ]
    
    for pattern in explicit_patterns:
        match = re.search(pattern, msg, re.IGNORECASE)
        if match:
            query = match.group(1).strip()
            if len(query) > 2:
                return True, query
    
    # 시사/최신 정보 키워드
    time_keywords = ['최신', '현재', '요즘', '지금', '오늘', '이번', '2024', '2025', '2026']
    info_keywords = ['뉴스', '소식', '가격', '환율', '주가', '시세', '날씨', '기온', '발표', '출시']
    
    for keyword in time_keywords + info_keywords:
        if keyword in msg:
            return True, msg
    
    return False, ""


@app.post("/chat")
async def chat_endpoint(chat_id: str = Form(...), message: str = Form(default=""), files: List[UploadFile] = File(default=[])):
    global chats
    
    if chat_id not in chats:
        chats[chat_id] = {"title": "새 채팅", "messages": [], "created": datetime.now().isoformat(), "updated": datetime.now().isoformat()}
    
    try:
        user_message = message.strip()
        file_contents = []
        file_names = []
        
        for file in files:
            if file.filename:
                try:
                    file_bytes = await file.read()
                    file_text = extract_file_content(file_bytes, file.filename)
                    if file_text:
                        file_contents.append(f"[파일: {file.filename}]\n{file_text[:25000]}")
                        file_names.append(file.filename)
                except: pass
        
        if file_contents:
            final_content = "\n\n".join(file_contents) + f"\n\n질문: {user_message or '위 문서를 분석해주세요.'}"
            display_content = user_message + (f" 📎 {', '.join(file_names)}" if user_message else f"📎 {', '.join(file_names)}")
        elif user_message:
            final_content = user_message
            display_content = user_message
            
            # 웹 검색 필요 여부 확인
            need_search, search_query = should_search(user_message)
            if need_search and search_query:
                search_results = await web_search(search_query)
                if search_results:
                    final_content = f"""[🔍 웹 검색 결과: "{search_query}"]

{search_results}

---
위 검색 결과를 참고하여 다음 질문에 답해주세요. 검색 결과의 정보를 활용하되, 출처를 명시해주세요.

질문: {user_message}"""
                    display_content = f"🔍 {user_message}"
        else:
            return JSONResponse({"response": "메시지를 입력해주세요.", "tokens_used": 0})
        
        # 메시지 저장 (표시용과 API용 분리)
        chats[chat_id]["messages"].append({"role": "user", "content": final_content, "display": display_content, "time": datetime.now().isoformat()})
        chats[chat_id]["updated"] = datetime.now().isoformat()
        
        # 첫 메시지면 제목 생성
        if len(chats[chat_id]["messages"]) == 1:
            chats[chat_id]["title"] = generate_title(user_message or file_names[0] if file_names else "PDF 분석")
        
        # API 호출용 메시지 (display 제외)
        api_messages = [{"role": m["role"], "content": m["content"]} for m in chats[chat_id]["messages"]]
        
        response = client.messages.create(
            model="claude-opus-4-20250514", max_tokens=6000, system=CACHED_SYSTEM,
            messages=api_messages, extra_headers={"anthropic-beta": "prompt-caching-2024-07-31"}
        )
        
        assistant_message = response.content[0].text
        chats[chat_id]["messages"].append({"role": "assistant", "content": assistant_message, "display": assistant_message, "time": datetime.now().isoformat()})
        save_chats()
        
        return JSONResponse({
            "response": assistant_message,
            "tokens_used": response.usage.input_tokens + response.usage.output_tokens,
            "title": chats[chat_id]["title"],
            "cache_read": getattr(response.usage, 'cache_read_input_tokens', 0),
            "cache_create": getattr(response.usage, 'cache_creation_input_tokens', 0)
        })
        
    except RateLimitError:
        if chats[chat_id]["messages"] and chats[chat_id]["messages"][-1]["role"] == "user":
            chats[chat_id]["messages"].pop()
        return JSONResponse({"response": "⚠️ API 요청 한도 초과. 잠시 후 다시 시도해주세요.", "tokens_used": 0})
    except APIConnectionError:
        if chats[chat_id]["messages"] and chats[chat_id]["messages"][-1]["role"] == "user":
            chats[chat_id]["messages"].pop()
        return JSONResponse({"response": "⚠️ 연결 오류. 인터넷 연결을 확인해주세요.", "tokens_used": 0})
    except APIStatusError as e:
        if chats[chat_id]["messages"] and chats[chat_id]["messages"][-1]["role"] == "user":
            chats[chat_id]["messages"].pop()
        return JSONResponse({"response": f"⚠️ API 오류: {e.message}", "tokens_used": 0})
    except Exception as e:
        if chats[chat_id]["messages"] and chats[chat_id]["messages"][-1]["role"] == "user":
            chats[chat_id]["messages"].pop()
        print(traceback.format_exc())
        return JSONResponse({"response": f"⚠️ 오류: {e}", "tokens_used": 0})

@app.get("/chats")
async def get_chats():
    return JSONResponse([
        {"id": k, "title": v["title"], "created": v["created"], "updated": v.get("updated", v["created"]), "messageCount": len(v["messages"])}
        for k, v in sorted(chats.items(), key=lambda x: x[1].get("updated", x[1]["created"]), reverse=True)
    ])

@app.get("/chat/{chat_id}")
async def get_chat(chat_id: str):
    if chat_id in chats:
        return JSONResponse(chats[chat_id])
    return JSONResponse({"messages": [], "title": "새 채팅"})

@app.delete("/chat/{chat_id}")
async def delete_chat(chat_id: str):
    if chat_id in chats:
        del chats[chat_id]
        save_chats()
    return JSONResponse({"status": "deleted"})

@app.put("/chat/{chat_id}/title")
async def update_title(chat_id: str, title: str = Form(...)):
    if chat_id in chats:
        chats[chat_id]["title"] = title
        save_chats()
    return JSONResponse({"status": "updated"})

@app.get("/search")
async def search_chats(q: str = Query(...)):
    results = []
    for chat_id, chat in chats.items():
        for i, msg in enumerate(chat["messages"]):
            if q.lower() in msg.get("display", msg["content"]).lower():
                results.append({
                    "chatId": chat_id,
                    "chatTitle": chat["title"],
                    "messageIndex": i,
                    "preview": msg.get("display", msg["content"])[:100],
                    "role": msg["role"]
                })
    return JSONResponse(results[:20])

@app.get("/export/{chat_id}")
async def export_chat(chat_id: str, format: str = "md"):
    if chat_id not in chats:
        return JSONResponse({"error": "채팅을 찾을 수 없습니다"}, status_code=404)
    
    chat = chats[chat_id]
    if format == "md":
        content = f"# {chat['title']}\n\n생성: {chat['created']}\n\n---\n\n"
        for msg in chat["messages"]:
            role = "👤 나" if msg["role"] == "user" else "🤖 AI"
            content += f"## {role}\n\n{msg.get('display', msg['content'])}\n\n---\n\n"
        return StreamingResponse(
            iter([content.encode('utf-8')]),
            media_type="text/markdown",
            headers={"Content-Disposition": f"attachment; filename={chat_id}.md"}
        )
    else:
        return JSONResponse(chat)

@app.get("/settings")
async def get_settings():
    return JSONResponse(settings)

@app.post("/settings")
async def update_settings(theme: str = Form(None), fontSize: str = Form(None)):
    if theme: settings["theme"] = theme
    if fontSize: settings["fontSize"] = fontSize
    save_settings()
    return JSONResponse(settings)

@app.get("/stats")
async def get_stats():
    total_messages = sum(len(c["messages"]) for c in chats.values())
    return JSONResponse({
        "totalChats": len(chats),
        "totalMessages": total_messages,
        "oldestChat": min((c["created"] for c in chats.values()), default=None),
        "newestChat": max((c.get("updated", c["created"]) for c in chats.values()), default=None)
    })

@app.get("/web-search")
async def web_search_endpoint(q: str = Query(...)):
    """수동 웹 검색 엔드포인트"""
    results = await web_search(q, num_results=8)
    return JSONResponse({"query": q, "results": results})

@app.get("/", response_class=HTMLResponse)
async def index():
    return HTML

HTML = '''<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>HarimCraft - 코딩 챗봇</title>
<link href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.5.0/css/all.min.css" rel="stylesheet">
<link href="https://cdnjs.cloudflare.com/ajax/libs/KaTeX/0.16.9/katex.min.css" rel="stylesheet">
<link href="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.9.0/styles/atom-one-dark.min.css" rel="stylesheet">
<style>
@import url('https://fonts.googleapis.com/css2?family=VT323&display=swap');
*{margin:0;padding:0;box-sizing:border-box}
:root{
    --bg:#2d2d2d;--bg2:#3d3d3d;--bg3:#4d4d4d;--bg4:#5d5d5d;
    --accent:#5b8731;--accent2:#7cb342;--accent-glow:rgba(91,135,49,.4);
    --text:#e8e8e8;--text2:#b0b0b0;--text3:#808080;
    --border:#1a1a1a;--success:#5b8731;--warning:#c6a000;--error:#b02e26;
    --dirt:#866043;--stone:#7d7d7d;--grass:#5b8731;--wood:#9c6d3e;
    --diamond:#4aedd9;--gold:#fcdb05;--redstone:#ff0000;
    --gradient:linear-gradient(135deg,#5b8731,#7cb342);
}
.light{
    --bg:#c6c6c6;--bg2:#d4d4d4;--bg3:#e2e2e2;--bg4:#f0f0f0;
    --text:#1a1a1a;--text2:#3d3d3d;--text3:#5d5d5d;--border:#a0a0a0;
}
body{font-family:'VT323','Pretendard',monospace;background:var(--bg);color:var(--text);height:100vh;display:flex;overflow:hidden;font-size:18px;image-rendering:pixelated}
body.font-small{font-size:16px}
body.font-large{font-size:22px}

/* 픽셀 보더 효과 */
.pixel-border{
    border:4px solid;
    border-color:#fff #555 #555 #fff;
    box-shadow:inset 2px 2px 0 rgba(255,255,255,.2),inset -2px -2px 0 rgba(0,0,0,.2);
}

/* 사이드바 */
.sidebar{width:280px;background:var(--bg2);border-right:4px solid var(--border);display:flex;flex-direction:column;flex-shrink:0;transition:transform .3s}
.sidebar.collapsed{transform:translateX(-280px);position:absolute;z-index:100;height:100%}
.sidebar-header{padding:1.25rem;border-bottom:4px solid var(--border)}
.logo{display:flex;align-items:center;gap:.75rem;margin-bottom:1rem}
.logo-icon{width:44px;height:44px;background:var(--grass);border:3px solid;border-color:#7cb342 #3d5c1f #3d5c1f #7cb342;display:flex;align-items:center;justify-content:center;font-size:1.5rem}
.logo-text{font-size:1.4rem;font-weight:700;letter-spacing:1px}
.logo-text span{color:var(--grass)}
.new-chat-btn{width:100%;padding:1rem;background:var(--grass);color:#fff;border:3px solid;border-color:#7cb342 #3d5c1f #3d5c1f #7cb342;cursor:pointer;font-size:1.1rem;font-family:inherit;display:flex;align-items:center;justify-content:center;gap:.5rem;transition:all .1s}
.new-chat-btn:hover{filter:brightness(1.1)}
.new-chat-btn:active{border-color:#3d5c1f #7cb342 #7cb342 #3d5c1f}

/* 검색 */
.search-box{padding:1rem;border-bottom:4px solid var(--border)}
.search-input{width:100%;padding:.75rem 1rem .75rem 2.5rem;background:var(--bg);border:3px solid;border-color:#555 #fff #fff #555;color:var(--text);font-size:1rem;font-family:inherit}
.search-input:focus{outline:none;border-color:var(--grass) #3d5c1f #3d5c1f var(--grass)}
.search-wrapper{position:relative}
.search-wrapper i{position:absolute;left:.875rem;top:50%;transform:translateY(-50%);color:var(--text3)}

/* 채팅 목록 */
.chat-list{flex:1;overflow-y:auto;padding:.5rem}
.chat-item{padding:.875rem 1rem;cursor:pointer;display:flex;align-items:center;gap:.75rem;margin-bottom:.25rem;transition:all .1s;border:2px solid transparent}
.chat-item:hover{background:var(--bg3);border-color:var(--border)}
.chat-item.active{background:var(--bg4);border:2px solid var(--grass)}
.chat-item-icon{width:32px;height:32px;background:var(--dirt);border:2px solid;border-color:#a07850 #5c4030 #5c4030 #a07850;display:flex;align-items:center;justify-content:center;color:#fff;font-size:.9rem}
.chat-item-content{flex:1;min-width:0}
.chat-item-title{font-size:1rem;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.chat-item-meta{font-size:.85rem;color:var(--text3);margin-top:.125rem}
.chat-item-actions{opacity:0;display:flex;gap:.25rem}
.chat-item:hover .chat-item-actions{opacity:1}
.chat-item-btn{padding:.375rem;border:0;background:0;color:var(--text3);cursor:pointer;font-size:1rem}
.chat-item-btn:hover{color:var(--text)}
.chat-item-btn.delete:hover{color:var(--error)}

/* 메인 영역 */
.main{flex:1;display:flex;flex-direction:column;min-width:0;background:var(--bg)}
.header{background:var(--bg2);padding:1rem 1.5rem;border-bottom:4px solid var(--border);display:flex;align-items:center;gap:1rem}
.menu-btn{display:none;padding:.5rem;border:0;background:0;color:var(--text2);cursor:pointer;font-size:1.2rem}
.header-title{flex:1;font-size:1.2rem;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.header-actions{display:flex;gap:.5rem}
.header-btn{padding:.5rem .75rem;border:3px solid;border-color:#555 #fff #fff #555;background:var(--bg3);color:var(--text2);cursor:pointer;font-size:1rem;font-family:inherit;display:flex;align-items:center;gap:.375rem;transition:all .1s}
.header-btn:hover{background:var(--bg4)}
.header-btn:active{border-color:#fff #555 #555 #fff}

/* 채팅 영역 */
#chat{flex:1;overflow-y:auto;padding:2rem;display:flex;flex-direction:column;gap:1.5rem;background:linear-gradient(var(--bg) 0%,var(--bg) 100%)}
.message{display:flex;gap:1rem;animation:fadeIn .3s;max-width:850px;width:100%;margin:0 auto}
@keyframes fadeIn{from{opacity:0;transform:translateY(10px)}to{opacity:1}}
.message.user{flex-direction:row-reverse}
.avatar{width:40px;height:40px;border:3px solid;display:flex;align-items:center;justify-content:center;font-size:1rem;flex-shrink:0}
.message.assistant .avatar{background:var(--grass);border-color:#7cb342 #3d5c1f #3d5c1f #7cb342}
.message.user .avatar{background:var(--diamond);border-color:#7fffff #2a9d9d #2a9d9d #7fffff}
.bubble{max-width:calc(100% - 50px);padding:1rem 1.25rem;line-height:1.6;border:3px solid}
.message.user .bubble{background:#3d7a9e;color:#fff;border-color:#5ba3c9 #2a5570 #2a5570 #5ba3c9}
.message.assistant .bubble{background:var(--bg2);border-color:#555 #1a1a1a #1a1a1a #555}
.bubble p{margin:.5rem 0}
.bubble ul,.bubble ol{margin:.5rem 0 .5rem 1.5rem}
.bubble li{margin:.25rem 0}
.bubble strong{color:var(--gold)}
.bubble a{color:var(--diamond);text-decoration:underline}

/* 코드 블록 */
.code-block{position:relative;margin:1rem 0;overflow:hidden;background:#1a1a1a;border:3px solid;border-color:#333 #000 #000 #333}
.code-header{display:flex;justify-content:space-between;align-items:center;padding:.625rem 1rem;background:#2d2d2d;font-size:1rem;color:#b0b0b0;border-bottom:2px solid #000}
.code-lang{display:flex;align-items:center;gap:.5rem}
.code-lang i{color:var(--grass)}
.copy-btn{background:var(--bg3);border:2px solid;border-color:#555 #222 #222 #555;color:#b0b0b0;padding:.375rem .75rem;cursor:pointer;font-size:1rem;font-family:inherit;display:flex;align-items:center;gap:.375rem;transition:all .1s}
.copy-btn:hover{background:var(--bg4)}
.copy-btn.copied{background:var(--grass);color:#fff}
.code-block pre{margin:0;padding:1rem;overflow-x:auto}
.code-block code{font-family:'VT323',monospace;font-size:1.1rem;line-height:1.5}
.bubble code:not(.hljs){background:#1a1a1a;padding:.2rem .5rem;font-size:1rem;color:var(--gold);font-family:'VT323',monospace;border:2px solid #333}

/* 입력 영역 */
.input-area{background:var(--bg2);border-top:4px solid var(--border);padding:1.25rem 2rem}
.input-wrapper{max-width:850px;margin:0 auto;display:flex;flex-direction:column;gap:.75rem}
.files-preview{display:none;flex-wrap:wrap;gap:.5rem}
.files-preview.active{display:flex}
.file-tag{display:flex;align-items:center;gap:.5rem;padding:.5rem .875rem;background:var(--bg3);font-size:1rem;color:var(--text);border:2px solid;border-color:#555 #222 #222 #555}
.file-tag i.fa-file-pdf{color:var(--error)}
.file-tag i.fa-file-word{color:#5ba3c9}
.file-tag i.fa-file-excel{color:var(--grass)}
.file-tag i.fa-file-powerpoint{color:var(--warning)}
.file-tag i.fa-file-code{color:var(--diamond)}
.file-tag i.fa-file-csv{color:var(--diamond)}
.file-tag .remove{cursor:pointer;color:var(--text3);margin-left:.25rem}
.file-tag .remove:hover{color:var(--error)}
.input-row{display:flex;gap:.75rem;align-items:flex-end}
.input-box{flex:1;display:flex;align-items:flex-end;background:var(--bg);border:3px solid;border-color:#555 #fff #fff #555;padding:.5rem;transition:all .1s}
.input-box:focus-within{border-color:var(--grass) #3d5c1f #3d5c1f var(--grass)}
#msgInput{flex:1;background:0;border:0;color:var(--text);font-size:1.1rem;font-family:inherit;padding:.75rem;resize:none;max-height:150px;outline:0;line-height:1.4}
#msgInput::placeholder{color:var(--text3)}
.input-actions{display:flex;gap:.25rem;padding:.25rem}
.icon-btn{width:40px;height:40px;border:2px solid;border-color:#555 #222 #222 #555;background:var(--bg3);color:var(--text3);cursor:pointer;display:flex;align-items:center;justify-content:center;transition:all .1s;font-size:1.1rem}
.icon-btn:hover{background:var(--bg4);color:var(--text)}
.send-btn{width:50px;height:50px;border:3px solid;border-color:#7cb342 #3d5c1f #3d5c1f #7cb342;background:var(--grass);color:#fff;cursor:pointer;display:flex;align-items:center;justify-content:center;font-size:1.2rem;transition:all .1s}
.send-btn:hover{filter:brightness(1.1)}
.send-btn:active{border-color:#3d5c1f #7cb342 #7cb342 #3d5c1f}
.send-btn:disabled{opacity:.5;cursor:not-allowed}

/* 타이핑 인디케이터 */
.typing-dots{display:flex;gap:6px;padding:.5rem}
.typing-dots span{width:10px;height:10px;background:var(--grass);animation:bounce 1.4s infinite}
.typing-dots span:nth-child(1){animation-delay:-.32s}
.typing-dots span:nth-child(2){animation-delay:-.16s}
@keyframes bounce{0%,80%,100%{transform:scale(0)}40%{transform:scale(1)}}

/* 토큰 정보 */
.token-info{display:flex;justify-content:center;gap:1rem;font-size:1rem;color:var(--text3);margin-top:.5rem}
.token-info span{display:flex;align-items:center;gap:.25rem}

/* 빈 상태 */
.empty-state{flex:1;display:flex;flex-direction:column;align-items:center;justify-content:center;color:var(--text3);text-align:center;padding:2rem}
.empty-icon{width:80px;height:80px;background:var(--grass);border:4px solid;border-color:#7cb342 #3d5c1f #3d5c1f #7cb342;display:flex;align-items:center;justify-content:center;font-size:2.5rem;margin-bottom:1.5rem}
.empty-title{font-size:1.5rem;color:var(--text);margin-bottom:.5rem}
.empty-desc{max-width:300px;line-height:1.5;font-size:1.1rem}

/* 설정 모달 */
.modal-overlay{position:fixed;inset:0;background:rgba(0,0,0,.7);display:none;align-items:center;justify-content:center;z-index:1000}
.modal-overlay.active{display:flex}
.modal{background:var(--bg2);padding:1.5rem;width:90%;max-width:400px;border:4px solid;border-color:#555 #1a1a1a #1a1a1a #555}
.modal-header{display:flex;justify-content:space-between;align-items:center;margin-bottom:1.5rem}
.modal-title{font-size:1.3rem}
.modal-close{padding:.5rem;border:2px solid;border-color:#555 #222 #222 #555;background:var(--bg3);color:var(--text2);cursor:pointer;font-family:inherit;font-size:1rem}
.modal-close:hover{background:var(--bg4)}
.setting-item{margin-bottom:1.25rem}
.setting-label{font-size:1rem;color:var(--text2);margin-bottom:.5rem}
.setting-options{display:flex;gap:.5rem}
.setting-btn{flex:1;padding:.75rem;border:3px solid;border-color:#555 #222 #222 #555;background:var(--bg3);color:var(--text);cursor:pointer;font-size:1rem;font-family:inherit;transition:all .1s}
.setting-btn:hover{background:var(--bg4)}
.setting-btn.active{background:var(--grass);border-color:#7cb342 #3d5c1f #3d5c1f #7cb342;color:#fff}

/* 검색 결과 */
.search-results{position:absolute;top:100%;left:0;right:0;background:var(--bg2);border:3px solid;border-color:#555 #1a1a1a #1a1a1a #555;margin-top:.5rem;max-height:300px;overflow-y:auto;z-index:10;display:none}
.search-results.active{display:block}
.search-result-item{padding:.75rem 1rem;cursor:pointer;border-bottom:2px solid var(--border)}
.search-result-item:hover{background:var(--bg3)}
.search-result-item:last-child{border-bottom:0}
.search-result-title{font-size:1rem;color:var(--grass);margin-bottom:.25rem}
.search-result-preview{font-size:.95rem;color:var(--text2);overflow:hidden;text-overflow:ellipsis;white-space:nowrap}

/* 스크롤바 */
::-webkit-scrollbar{width:10px}
::-webkit-scrollbar-track{background:var(--bg)}
::-webkit-scrollbar-thumb{background:var(--stone);border:2px solid var(--bg)}

/* 반응형 */
@media(max-width:768px){
    .sidebar{position:absolute;z-index:100;height:100%;transform:translateX(-280px)}
    .sidebar.open{transform:translateX(0)}
    .menu-btn{display:flex}
}
</style>
</head>
<body>
<aside class="sidebar" id="sidebar">
<div class="sidebar-header">
<div class="logo">
<div class="logo-icon">⛏️</div>
<div class="logo-text"><span>Harim</span>Craft</div>
</div>
<button class="new-chat-btn" onclick="newChat()"><i class="fas fa-plus"></i> 새 채팅</button>
</div>
<div class="search-box">
<div class="search-wrapper">
<i class="fas fa-search"></i>
<input type="text" class="search-input" id="searchInput" placeholder="채팅 검색..." oninput="searchChats(this.value)">
<div class="search-results" id="searchResults"></div>
</div>
</div>
<div class="chat-list" id="chatList"></div>
</aside>

<main class="main">
<header class="header">
<button class="menu-btn" onclick="toggleSidebar()"><i class="fas fa-bars"></i></button>
<div class="header-title" id="headerTitle">새 채팅</div>
<div class="header-actions">
<button class="header-btn" onclick="exportChat()" title="내보내기"><i class="fas fa-download"></i></button>
<button class="header-btn" onclick="openSettings()" title="설정"><i class="fas fa-cog"></i></button>
</div>
</header>

<div id="chat">
<div class="empty-state">
<div class="empty-icon"><i class="fas fa-comments"></i></div>
<div class="empty-title">안녕하세요, 하림님!</div>
<div class="empty-desc">퀀트 연구, 논문 리뷰, Python 코딩 무엇이든 물어보세요.</div>
</div>
</div>

<div class="input-area">
<div class="input-wrapper">
<div class="files-preview" id="filesPreview"></div>
<div class="input-row">
<div class="input-box">
<textarea id="msgInput" placeholder="메시지를 입력하세요..." rows="1"></textarea>
<div class="input-actions">
<input type="file" id="fileInput" accept=".pdf,.docx,.xlsx,.xls,.pptx,.csv,.txt,.md,.py,.js,.ts,.java,.json,.xml,.html,.css,.sql,.yaml,.yml" multiple hidden>
<button class="icon-btn" onclick="document.getElementById('fileInput').click()" title="PDF 첨부"><i class="fas fa-paperclip"></i></button>
</div>
</div>
<button class="send-btn" id="sendBtn" onclick="sendMessage()" title="전송"><i class="fas fa-paper-plane"></i></button>
</div>
<div class="token-info" id="tokenInfo"></div>
</div>
</div>
</main>

<!-- 설정 모달 -->
<div class="modal-overlay" id="settingsModal">
<div class="modal">
<div class="modal-header">
<div class="modal-title">설정</div>
<button class="modal-close" onclick="closeSettings()"><i class="fas fa-times"></i></button>
</div>
<div class="setting-item">
<div class="setting-label">테마</div>
<div class="setting-options">
<button class="setting-btn" data-theme="dark" onclick="setTheme('dark')"><i class="fas fa-moon"></i> 다크</button>
<button class="setting-btn" data-theme="light" onclick="setTheme('light')"><i class="fas fa-sun"></i> 라이트</button>
</div>
</div>
<div class="setting-item">
<div class="setting-label">글자 크기</div>
<div class="setting-options">
<button class="setting-btn" data-font="small" onclick="setFontSize('small')">작게</button>
<button class="setting-btn" data-font="medium" onclick="setFontSize('medium')">보통</button>
<button class="setting-btn" data-font="large" onclick="setFontSize('large')">크게</button>
</div>
</div>
</div>
</div>

<script src="https://cdnjs.cloudflare.com/ajax/libs/marked/11.1.1/marked.min.js"></script>
<script src="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.9.0/highlight.min.js"></script>
<script src="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.9.0/languages/python.min.js"></script>
<script src="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.9.0/languages/javascript.min.js"></script>
<script src="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.9.0/languages/sql.min.js"></script>
<script src="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.9.0/languages/bash.min.js"></script>
<script src="https://cdnjs.cloudflare.com/ajax/libs/KaTeX/0.16.9/katex.min.js"></script>
<script src="https://cdnjs.cloudflare.com/ajax/libs/KaTeX/0.16.9/contrib/auto-render.min.js"></script>
'''
HTML += '''
<script>
let currentChatId = null;
let selectedFiles = [];
let settings = {theme: 'dark', fontSize: 'medium'};

const chat = document.getElementById('chat');
const msgInput = document.getElementById('msgInput');
const sendBtn = document.getElementById('sendBtn');
const fileInput = document.getElementById('fileInput');
const filesPreview = document.getElementById('filesPreview');
const tokenInfo = document.getElementById('tokenInfo');
const chatList = document.getElementById('chatList');
const headerTitle = document.getElementById('headerTitle');
const searchInput = document.getElementById('searchInput');
const searchResults = document.getElementById('searchResults');
const sidebar = document.getElementById('sidebar');

// Marked 설정
const renderer = new marked.Renderer();
renderer.code = function(code, lang) {
    const language = lang || 'plaintext';
    let highlighted;
    try {
        highlighted = lang && hljs.getLanguage(lang) ? hljs.highlight(code, {language}).value : hljs.highlightAuto(code).value;
    } catch(e) {
        highlighted = code;
    }
    const id = 'code-' + Math.random().toString(36).substr(2,9);
    const langIcon = {'python':'fab fa-python','javascript':'fab fa-js','sql':'fas fa-database','bash':'fas fa-terminal'}[language] || 'fas fa-code';
    return `<div class="code-block"><div class="code-header"><span class="code-lang"><i class="${langIcon}"></i> ${language}</span><button class="copy-btn" onclick="copyCode('${id}')"><i class="fas fa-copy"></i> 복사</button></div><pre><code id="${id}" class="hljs">${highlighted}</code></pre></div>`;
};
marked.use({renderer, breaks: true});

function copyCode(id) {
    const code = document.getElementById(id);
    navigator.clipboard.writeText(code.textContent).then(() => {
        const btn = code.closest('.code-block').querySelector('.copy-btn');
        btn.innerHTML = '<i class="fas fa-check"></i> 복사됨';
        btn.classList.add('copied');
        setTimeout(() => { btn.innerHTML = '<i class="fas fa-copy"></i> 복사'; btn.classList.remove('copied'); }, 2000);
    });
}

// 유틸리티
function generateId() { return Date.now().toString(36) + Math.random().toString(36).substr(2); }
function formatDate(iso) {
    const d = new Date(iso);
    const now = new Date();
    const diff = now - d;
    if (diff < 60000) return '방금 전';
    if (diff < 3600000) return Math.floor(diff/60000) + '분 전';
    if (diff < 86400000) return Math.floor(diff/3600000) + '시간 전';
    if (diff < 604800000) return Math.floor(diff/86400000) + '일 전';
    return d.toLocaleDateString('ko-KR');
}

// 설정
async function loadSettings() {
    try {
        const res = await fetch('/settings');
        settings = await res.json();
        applySettings();
    } catch(e) {}
}

function applySettings() {
    document.body.classList.toggle('light', settings.theme === 'light');
    document.body.classList.remove('font-small', 'font-large');
    if (settings.fontSize !== 'medium') document.body.classList.add('font-' + settings.fontSize);
    document.querySelectorAll('[data-theme]').forEach(b => b.classList.toggle('active', b.dataset.theme === settings.theme));
    document.querySelectorAll('[data-font]').forEach(b => b.classList.toggle('active', b.dataset.font === settings.fontSize));
}

async function setTheme(theme) {
    settings.theme = theme;
    applySettings();
    await fetch('/settings', {method:'POST', body: new URLSearchParams({theme})});
}

async function setFontSize(size) {
    settings.fontSize = size;
    applySettings();
    await fetch('/settings', {method:'POST', body: new URLSearchParams({fontSize: size})});
}

function openSettings() { document.getElementById('settingsModal').classList.add('active'); }
function closeSettings() { document.getElementById('settingsModal').classList.remove('active'); }
function toggleSidebar() { sidebar.classList.toggle('open'); }

// 채팅 목록
async function loadChatList() {
    try {
        const res = await fetch('/chats');
        const chats = await res.json();
        chatList.innerHTML = chats.map(c => `
            <div class="chat-item ${c.id === currentChatId ? 'active' : ''}" onclick="loadChat('${c.id}')">
                <div class="chat-item-icon"><i class="fas fa-message"></i></div>
                <div class="chat-item-content">
                    <div class="chat-item-title">${escapeHtml(c.title)}</div>
                    <div class="chat-item-meta">${formatDate(c.updated)} · ${c.messageCount}개 메시지</div>
                </div>
                <div class="chat-item-actions">
                    <button class="chat-item-btn" onclick="event.stopPropagation();renameChat('${c.id}','${escapeHtml(c.title)}')" title="이름 변경"><i class="fas fa-pen"></i></button>
                    <button class="chat-item-btn delete" onclick="event.stopPropagation();deleteChat('${c.id}')" title="삭제"><i class="fas fa-trash"></i></button>
                </div>
            </div>
        `).join('') || '<div style="padding:1rem;text-align:center;color:var(--text3)">채팅이 없습니다</div>';
    } catch(e) {}
}

function escapeHtml(str) {
    return str.replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;');
}

async function loadChat(chatId) {
    currentChatId = chatId;
    try {
        const res = await fetch('/chat/' + chatId);
        const data = await res.json();
        headerTitle.textContent = data.title || '새 채팅';
        chat.innerHTML = '';
        data.messages.forEach(m => addMsg(m.display || m.content, m.role === 'user', false));
        if (!data.messages.length) showEmptyState();
        loadChatList();
        chat.scrollTop = chat.scrollHeight;
        sidebar.classList.remove('open');
    } catch(e) {}
}

function newChat() {
    currentChatId = generateId();
    headerTitle.textContent = '새 채팅';
    showEmptyState();
    tokenInfo.innerHTML = '';
    loadChatList();
    sidebar.classList.remove('open');
}

function showEmptyState() {
    chat.innerHTML = `<div class="empty-state">
        <div class="empty-icon">⛏️</div>
        <div class="empty-title">HarimCraft에 오신 걸 환영해요!</div>
        <div class="empty-desc">퀀트 연구, 논문 리뷰, Python 코딩 무엇이든 물어보세요.</div>
    </div>`;
}

async function deleteChat(chatId) {
    if (!confirm('이 채팅을 삭제할까요?')) return;
    await fetch('/chat/' + chatId, {method: 'DELETE'});
    if (chatId === currentChatId) newChat();
    else loadChatList();
}

async function renameChat(chatId, currentTitle) {
    const newTitle = prompt('새 이름을 입력하세요:', currentTitle);
    if (newTitle && newTitle !== currentTitle) {
        await fetch('/chat/' + chatId + '/title', {method:'PUT', body: new URLSearchParams({title: newTitle})});
        if (chatId === currentChatId) headerTitle.textContent = newTitle;
        loadChatList();
    }
}

async function exportChat() {
    if (!currentChatId) return;
    window.open('/export/' + currentChatId + '?format=md', '_blank');
}

// 검색
let searchTimeout;
async function searchChats(query) {
    clearTimeout(searchTimeout);
    if (!query.trim()) {
        searchResults.classList.remove('active');
        return;
    }
    searchTimeout = setTimeout(async () => {
        try {
            const res = await fetch('/search?q=' + encodeURIComponent(query));
            const results = await res.json();
            if (results.length) {
                searchResults.innerHTML = results.map(r => `
                    <div class="search-result-item" onclick="loadChat('${r.chatId}');searchResults.classList.remove('active');searchInput.value='';">
                        <div class="search-result-title">${escapeHtml(r.chatTitle)}</div>
                        <div class="search-result-preview">${escapeHtml(r.preview)}</div>
                    </div>
                `).join('');
                searchResults.classList.add('active');
            } else {
                searchResults.innerHTML = '<div style="padding:1rem;text-align:center;color:var(--text3)">결과 없음</div>';
                searchResults.classList.add('active');
            }
        } catch(e) {}
    }, 300);
}

// 파일 처리
msgInput.addEventListener('input', function() { this.style.height = 'auto'; this.style.height = Math.min(this.scrollHeight, 150) + 'px'; });
msgInput.addEventListener('keydown', function(e) { if (e.key === 'Enter' && !e.shiftKey) { e.preventDefault(); sendMessage(); } });
fileInput.addEventListener('change', function() { Array.from(this.files).forEach(f => { if (!selectedFiles.find(x => x.name === f.name)) selectedFiles.push(f); }); updateFiles(); });

function getFileIcon(filename) {
    const ext = filename.split('.').pop().toLowerCase();
    const icons = {
        'pdf': 'fa-file-pdf',
        'docx': 'fa-file-word', 'doc': 'fa-file-word',
        'xlsx': 'fa-file-excel', 'xls': 'fa-file-excel', 'csv': 'fa-file-csv',
        'pptx': 'fa-file-powerpoint', 'ppt': 'fa-file-powerpoint',
        'py': 'fa-file-code', 'js': 'fa-file-code', 'ts': 'fa-file-code', 
        'java': 'fa-file-code', 'json': 'fa-file-code', 'html': 'fa-file-code',
        'css': 'fa-file-code', 'sql': 'fa-file-code', 'xml': 'fa-file-code'
    };
    return icons[ext] || 'fa-file';
}

function updateFiles() {
    filesPreview.innerHTML = '';
    if (!selectedFiles.length) { filesPreview.classList.remove('active'); return; }
    filesPreview.classList.add('active');
    selectedFiles.forEach((f, i) => {
        const tag = document.createElement('div');
        tag.className = 'file-tag';
        const icon = getFileIcon(f.name);
        tag.innerHTML = `<i class="fas ${icon}"></i> ${escapeHtml(f.name)} <i class="fas fa-times remove" onclick="removeFile(${i})"></i>`;
        filesPreview.appendChild(tag);
    });
}
function removeFile(i) { selectedFiles.splice(i, 1); updateFiles(); }
function clearFiles() { selectedFiles = []; fileInput.value = ''; updateFiles(); }

// 메시지
function addMsg(content, isUser, scroll=true) {
    const empty = chat.querySelector('.empty-state');
    if (empty) empty.remove();
    
    const div = document.createElement('div');
    div.className = 'message ' + (isUser ? 'user' : 'assistant');
    div.innerHTML = `<div class="avatar"><i class="fas fa-${isUser ? 'user' : 'robot'}"></i></div><div class="bubble"></div>`;
    const bubble = div.querySelector('.bubble');
    
    if (isUser) {
        bubble.textContent = content;
    } else {
        try {
            bubble.innerHTML = marked.parse(content || '');
            renderMathInElement(bubble, {
                delimiters: [{left:'$$',right:'$$',display:true},{left:'$',right:'$',display:false},{left:'\\\\[',right:'\\\\]',display:true},{left:'\\\\(',right:'\\\\)',display:false}],
                throwOnError: false
            });
        } catch(e) { bubble.textContent = content || ''; }
    }
    chat.appendChild(div);
    if (scroll) chat.scrollTop = chat.scrollHeight;
}

function showTyping() {
    const div = document.createElement('div');
    div.className = 'message assistant';
    div.id = 'typing';
    div.innerHTML = '<div class="avatar"><i class="fas fa-robot"></i></div><div class="bubble"><div class="typing-dots"><span></span><span></span><span></span></div></div>';
    chat.appendChild(div);
    chat.scrollTop = chat.scrollHeight;
}
function hideTyping() { const t = document.getElementById('typing'); if(t) t.remove(); }

async function sendMessage() {
    const msg = msgInput.value.trim();
    if (!msg && !selectedFiles.length) return;
    if (!currentChatId) currentChatId = generateId();
    
    const displayMsg = msg + (selectedFiles.length ? ' 📎 ' + selectedFiles.map(f=>f.name).join(', ') : '');
    addMsg(displayMsg, true);
    msgInput.value = '';
    msgInput.style.height = 'auto';
    sendBtn.disabled = true;
    showTyping();
    
    const formData = new FormData();
    formData.append('chat_id', currentChatId);
    formData.append('message', msg);
    selectedFiles.forEach(f => formData.append('files', f));
    
    try {
        const ctrl = new AbortController();
        const timeout = setTimeout(() => ctrl.abort(), 180000);
        const res = await fetch('/chat', {method:'POST', body:formData, signal:ctrl.signal});
        clearTimeout(timeout);
        const data = await res.json();
        hideTyping();
        addMsg(data.response || '응답을 받지 못했습니다.', false);
        
        if (data.tokens_used > 0) {
            let info = `<span><i class="fas fa-coins"></i> ${data.tokens_used.toLocaleString()} 토큰</span>`;
            if (data.cache_read > 0) info += `<span><i class="fas fa-bolt"></i> 캐시 ${data.cache_read.toLocaleString()}</span>`;
            tokenInfo.innerHTML = info;
        }
        if (data.title) headerTitle.textContent = data.title;
        loadChatList();
    } catch(e) {
        hideTyping();
        addMsg('⚠️ ' + (e.name === 'AbortError' ? '요청 시간이 초과되었습니다.' : e.message), false);
    }
    clearFiles();
    sendBtn.disabled = false;
    msgInput.focus();
}

// 초기화
loadSettings();
loadChatList();
newChat();

// 클릭 외부 검색 닫기
document.addEventListener('click', e => {
    if (!e.target.closest('.search-wrapper')) searchResults.classList.remove('active');
});
</script>
</body>
</html>'''

if __name__ == "__main__":
    import uvicorn
    print("=" * 50)
    print("🚀 하림님 코딩 챗봇 v2.0")
    print("=" * 50)
    print("📍 http://localhost:8000")
    print("=" * 50)
    uvicorn.run(app, host="0.0.0.0", port=8000)
